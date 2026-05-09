from __future__ import annotations

import unicodedata
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Any
from zipfile import ZipFile

SLIDE_WIDE = (13.333, 7.5)
ACCENT = "0B6E69"
CHARCOAL = "172126"
MUTED = "5E6B73"
SERIES_COLORS = ["0B6E69", "C45A28", "4267AC", "7B6D8D", "717C36", "A13D63"]

TEXT_REPLACEMENTS = {
    "\u2018": "'",
    "\u2019": "'",
    "\u201c": '"',
    "\u201d": '"',
    "\u2013": "-",
    "\u2014": "-",
    "\u00a0": " ",
    "\u2026": "...",
}


@dataclass
class Series:
    label: str
    points: list[tuple[str, float | int]]


@dataclass
class ReportSlide:
    title: str
    subtitle: str
    chart_note: str
    notes: str
    source: str
    source_url: str
    unit_label: str
    series: list[Series]
    chart_type: str = "line"
    max_category_label_chars: int = 26


def inspect_pptx_report(path: str | Path) -> dict[str, Any]:
    """Return package-level checks for generated reports.

    Native editable charts appear as `ppt/charts/chart*.xml` with embedded
    workbooks. Report runs should not contain `ppt/media/*` chart images.
    """
    with ZipFile(Path(path)) as archive:
        names = archive.namelist()
    charts = [name for name in names if name.startswith("ppt/charts/chart") and name.endswith(".xml")]
    embedded_workbooks = [name for name in names if name.startswith("ppt/embeddings/")]
    media_files = [name for name in names if name.startswith("ppt/media/")]
    slides = [name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
    return {
        "slides": len(slides),
        "charts": len(charts),
        "embedded_workbooks": len(embedded_workbooks),
        "media_files": len(media_files),
        "native_chart_package": bool(charts) and len(charts) == len(embedded_workbooks) and not media_files,
    }


def write_pptx_report(
    title: str,
    question: str,
    executive_summary: str,
    slides: list[ReportSlide],
    output_path: str | Path,
) -> Path:
    """Write a PowerPoint deck with one native editable chart per evidence slide."""
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        raise RuntimeError("Install python-pptx to generate PowerPoint reports: pip install python-pptx") from exc

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDE[0])
    prs.slide_height = Inches(SLIDE_WIDE[1])
    blank = prs.slide_layouts[6]

    _validate_report(slides)
    _add_summary_slide(prs, blank, title, question, executive_summary)
    for index, report_slide in enumerate(slides, start=1):
        _add_chart_slide(prs, blank, index, report_slide)

    prs.save(path)
    return path


def _add_summary_slide(prs, layout, title: str, question: str, executive_summary: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(layout)
    _add_text(slide, "Pacific Data Hub Evidence Brief", 0.72, 0.42, 6.8, 0.28, Pt(12), ACCENT, bold=True)
    _add_text(slide, title, 0.72, 0.82, 11.8, 0.78, Pt(32), CHARCOAL, bold=True)
    _add_text(slide, f"Question: {question}\nDate: {date.today().isoformat()}", 0.72, 1.76, 11.4, 0.58, Pt(14), CHARCOAL)
    _add_text(slide, "Executive Summary", 0.72, 2.62, 5.2, 0.36, Pt(22), CHARCOAL, bold=True)
    body = slide.shapes.add_textbox(Inches(0.72), Inches(3.05), Inches(11.5), Inches(2.6)).text_frame
    body.word_wrap = True
    p = body.paragraphs[0]
    p.text = _safe_text(executive_summary)
    p.font.size = Pt(19)
    p.font.color.rgb = RGBColor.from_string(CHARCOAL)
    p.line_spacing = 1.08


def _add_chart_slide(prs, layout, index: int, panel: ReportSlide) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(layout)
    _add_text(slide, panel.title, 0.62, 0.38, 12.0, 0.45, Pt(22), CHARCOAL, bold=True)
    _add_text(slide, panel.subtitle, 0.62, 0.88, 12.0, 0.72, Pt(14), CHARCOAL)
    _add_text(slide, f"Figure {index}. {panel.chart_note}", 0.62, 1.58, 12.0, 0.28, Pt(9), MUTED)

    chart = _add_native_chart(slide, panel, 0.72, 1.98, 11.85, 3.75)
    _style_chart(chart, panel)

    if panel.notes:
        _add_text(slide, f"Notes: {panel.notes}", 0.72, 5.9, 11.7, 0.38, Pt(8.5), MUTED)
    _add_text(slide, f"Source: {panel.source} {panel.source_url}".strip(), 0.72, 6.42, 11.7, 0.28, Pt(8), MUTED)


def _add_native_chart(slide, panel: ReportSlide, left: float, top: float, width: float, height: float):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    categories = _ordered_categories(panel.series)
    chart_data = CategoryChartData()
    chart_data.categories = [_category_label(category, panel.max_category_label_chars) for category in categories]
    for series in panel.series:
        values_by_period = {period: value for period, value in series.points}
        chart_data.add_series(_safe_text(series.label), [values_by_period.get(period) for period in categories])

    chart_type = XL_CHART_TYPE.LINE_MARKERS
    if panel.chart_type.lower() == "bar":
        chart_type = XL_CHART_TYPE.BAR_CLUSTERED
    elif panel.chart_type.lower() == "column":
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    elif panel.chart_type.lower() == "stacked_bar":
        chart_type = XL_CHART_TYPE.BAR_STACKED
    elif panel.chart_type.lower() == "stacked_column":
        chart_type = XL_CHART_TYPE.COLUMN_STACKED

    frame = slide.shapes.add_chart(
        chart_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    )
    return frame.chart


def _style_chart(chart, panel: ReportSlide) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_LEGEND_POSITION
    from pptx.util import Pt

    chart.has_legend = len(panel.series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(8)
    category_axis.tick_labels.font.color.rgb = RGBColor.from_string(CHARCOAL)

    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(8)
    value_axis.tick_labels.font.color.rgb = RGBColor.from_string(CHARCOAL)
    value_axis.has_major_gridlines = True
    value_axis.axis_title.text_frame.text = panel.unit_label
    value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
    value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(MUTED)

    for index, series in enumerate(chart.series):
        fill = RGBColor.from_string(SERIES_COLORS[index % len(SERIES_COLORS)])
        series.format.line.color.rgb = fill
        series.format.line.width = Pt(2)
        try:
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = fill
            series.marker.format.line.color.rgb = fill
            series.marker.size = 5
        except Exception:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = fill


def _ordered_categories(series: list[Series]) -> list[str]:
    categories: list[str] = []
    seen: set[str] = set()
    for item in series:
        for period, _value in item.points:
            if period not in seen:
                seen.add(period)
                categories.append(period)
    return categories


def _validate_report(slides: list[ReportSlide]) -> None:
    if not slides:
        raise ValueError("PowerPoint reports need at least one chart slide.")
    for index, slide in enumerate(slides, start=1):
        _validate_slide(index, slide)


def _validate_slide(index: int, slide: ReportSlide) -> None:
    if not slide.series:
        raise ValueError(f"Slide {index} has no chart series.")
    categories = _ordered_categories(slide.series)
    if not categories:
        raise ValueError(f"Slide {index} has no chart categories.")
    category_set = set(categories)
    for series in slide.series:
        seen = set()
        for category, value in series.points:
            if category in seen:
                raise ValueError(f"Slide {index} series {series.label!r} repeats category {category!r}.")
            seen.add(category)
            if value is None:
                raise ValueError(f"Slide {index} series {series.label!r} has a missing value for {category!r}.")
            try:
                float(value)
            except (TypeError, ValueError) as exc:
                raise ValueError(f"Slide {index} series {series.label!r} has a non-numeric value for {category!r}.") from exc
        missing = category_set.difference(seen)
        if len(slide.series) > 1 and missing:
            missing_list = ", ".join(sorted(missing)[:4])
            raise ValueError(
                f"Slide {index} multi-series chart series {series.label!r} is missing categories: {missing_list}. "
                "Align categories before charting so PowerPoint does not render gaps as zero."
            )
    if slide.chart_type.lower() == "line":
        for series in slide.series:
            missing = category_set.difference({category for category, _value in series.points})
            if missing:
                missing_list = ", ".join(sorted(missing)[:4])
                raise ValueError(
                    f"Slide {index} line chart series {series.label!r} is missing categories: {missing_list}. "
                    "Use a complete line series, split the chart, or use a column comparison so gaps are not rendered as zero."
                )


def _safe_text(text: object) -> str:
    safe = str(text or "")
    for bad, good in TEXT_REPLACEMENTS.items():
        safe = safe.replace(bad, good)
    safe = unicodedata.normalize("NFKD", safe)
    return safe.encode("ascii", "ignore").decode("ascii")


def _category_label(category: str, max_chars: int) -> str:
    label = _safe_text(category)
    if max_chars > 3 and len(label) > max_chars:
        return f"{label[: max_chars - 3].rstrip()}..."
    return label


def _add_text(slide, text: str, left: float, top: float, width: float, height: float, font_size, color: str, bold: bool = False) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = _safe_text(text)
    paragraph.font.size = font_size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)
